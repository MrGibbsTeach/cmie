from __future__ import annotations

import copy
import json
import re
import textwrap
from pathlib import Path
from typing import Dict, Any, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .ai_lesson_engine import BASE_OUTPUT_DIR, slugify

TEMPLATE_PATH = Path(__file__).resolve().parent.parent / "assets" / "CMIE_Master_Template.pptx"

# --------------------------------------------------------------------
# Slide Design System
# --------------------------------------------------------------------

BG = RGBColor(246, 245, 242)          # warm off-white
PRIMARY = RGBColor(108, 99, 255)      # purple
SECONDARY = RGBColor(175, 169, 255)   # lavender
ACCENT = RGBColor(94, 208, 193)       # mint
TEXT = RGBColor(40, 40, 40)
MUTED = RGBColor(120, 120, 120)
CARD_BG = RGBColor(255, 255, 255)
LINE = RGBColor(220, 220, 220)


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_slide_bg(slide, color=BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_left_bar(slide) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.0), Inches(0.0), Inches(0.12), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()


def _add_blob(slide, left, top, width, height, color, transparency=25):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.fill.background()
    return shape


def _add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    return shape


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    font_size=20,
    bold=False,
    color=TEXT,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return box


def _add_bullet_block(
    slide,
    left,
    top,
    width,
    height,
    lines: List[str],
    font_size=20,
    color=TEXT,
    font_name="Aptos",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True

    if not lines:
        return box

    first = tf.paragraphs[0]
    first.text = f"• {lines[0]}"
    first.font.size = Pt(font_size)
    first.font.color.rgb = color
    first.font.name = font_name

    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name

    return box


def _prepare_title_lines(title: str, max_chars_per_line: int = 28, max_lines: int = 2) -> str:
    cleaned = " ".join((title or "").split())
    if not cleaned:
        return ""

    lines = textwrap.wrap(
        cleaned,
        width=max_chars_per_line,
        break_long_words=False,
        break_on_hyphens=False,
    )

    if len(lines) > max_lines:
        lines = lines[:max_lines]
        last = lines[-1].rstrip()
        if len(last) > max_chars_per_line - 3:
            last = last[: max_chars_per_line - 3].rstrip()
        lines[-1] = last.rstrip(" ,.-") + "..."

    return "\n".join(lines)


def _add_main_title(slide, title: str, top=0.45, font_size=28) -> None:
    wrapped = _prepare_title_lines(title, max_chars_per_line=28, max_lines=2)
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(top),
        Inches(8.3),
        Inches(1.0),
        text=wrapped,
        font_size=font_size,
        bold=True,
        color=TEXT,
        font_name="Aptos Display",
    )


def _add_subtitle(slide, subtitle: str) -> None:
    _add_textbox(
        slide,
        Inches(0.58),
        Inches(1.18),
        Inches(4.8),
        Inches(0.45),
        text=subtitle,
        font_size=10,
        bold=False,
        color=MUTED,
        font_name="Aptos",
    )


def _add_footer_tag(slide, text: str) -> None:
    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.55), Inches(6.55), Inches(1.45), Inches(0.45)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = CARD_BG
    tag.line.color.rgb = LINE
    tag.line.width = Pt(1)

    tf = tag.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT
    p.font.name = "Aptos"


# --------------------------------------------------------------------
# IO
# --------------------------------------------------------------------


def _load_lesson(path: Path) -> Dict[str, Any]:
    with path.open(encoding="utf-8") as f:
        return json.load(f)


# --------------------------------------------------------------------
# Template blueprint handling
# --------------------------------------------------------------------


def _normalize_title(text: str) -> str:
    text = (text or "").strip()
    text = text.replace("🔎", "").replace("🧠", "").replace("🌍", "").replace("🛠", "").replace("💬", "")
    text = " ".join(text.split())
    return text.lower()


def _slide_all_text(slide) -> str:
    out: list[str] = []
    for shp in slide.shapes:
        if hasattr(shp, "text_frame") and shp.has_text_frame:
            t = (shp.text_frame.text or "").strip()
            if t:
                out.append(t)
    return "\n".join(out)


def _index_blueprints(prs: Presentation) -> Dict[str, Any]:
    """
    Find blueprint slides in the template by their visible title text.
    """
    want = {
        "title": ["cmie master template"],
        "section": ["section title"],
        "objectives": ["lesson objectives"],
        "essential": ["essential question"],
        "explore": ["explore"],
        "hook": ["hook scenario"],
        "watch": ["watch & think", "watch and think"],
        "content": ["content slide title"],
        "real_world": ["real world example"],
        "activity": ["activity"],
        "reflection": ["reflection questions"],
        "teacher_notes": ["teacher notes"],
    }

    found: Dict[str, Any] = {}

    for slide in prs.slides:
        text = _normalize_title(_slide_all_text(slide))
        for key, aliases in want.items():
            if key in found:
                continue
            if any(alias == text or alias in text for alias in aliases):
                found[key] = slide

    missing = [k for k in want if k not in found]
    if missing:
        raise RuntimeError(
            "Template is missing blueprint slides (or their titles are different). "
            f"Missing keys: {missing}."
        )

    return found


def _duplicate_slide(prs: Presentation, source_slide):
    """
    Duplicate a slide by copying its XML. This preserves the full visual styling
    of the blueprint slide.
    """
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(blank_layout)

    # Remove any default placeholders from the new blank slide
    for shp in list(new_slide.shapes):
        if shp.is_placeholder:
            el = shp._element
            el.getparent().remove(el)

    # Copy shapes from blueprint slide
    for shp in source_slide.shapes:
        new_el = copy.deepcopy(shp._element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    return new_slide


def _delete_slides_by_rids(prs: Presentation, rids: set[str]) -> None:
    """
    Remove slides using their original relationship IDs.
    """
    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    for sldId in list(sldIdLst):
        if sldId.rId in rids:
            prs.part.drop_rel(sldId.rId)
            sldIdLst.remove(sldId)


# --------------------------------------------------------------------
# Text placement helpers
# --------------------------------------------------------------------


def _find_text_frame(slide, role: str):
    role_upper = role.strip().upper()

    # 1) Named shapes are best
    for shp in slide.shapes:
        try:
            if (shp.name or "").strip().upper() == role_upper and shp.has_text_frame:
                return shp.text_frame
        except Exception:
            continue

    # 2) Built-in title placeholder
    if role_upper == "TITLE":
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                return slide.shapes.title.text_frame
        except Exception:
            pass

    # 3) Placeholder fallback
    preferred = {
        "SUBTITLE": {PP_PLACEHOLDER.SUBTITLE},
        "BODY": {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT},
    }.get(role_upper, set())

    if preferred:
        for ph in getattr(slide, "placeholders", []):
            try:
                if ph.placeholder_format.type in preferred and ph.has_text_frame:
                    return ph.text_frame
            except Exception:
                continue

    # 4) Any text frame
    for shp in slide.shapes:
        try:
            if shp.has_text_frame:
                return shp.text_frame
        except Exception:
            continue

    return None


def _prepare_slide_title(title: str, max_chars_per_line: int = 32, max_lines: int = 2) -> str:
    """
    Wrap and shorten slide titles so they do not crash into body text.
    """
    cleaned = " ".join((title or "").split())
    if not cleaned:
        return ""

    lines = textwrap.wrap(
        cleaned,
        width=max_chars_per_line,
        break_long_words=False,
        break_on_hyphens=False,
    )

    if not lines:
        return cleaned

    if len(lines) > max_lines:
        lines = lines[:max_lines]
        last = lines[-1].rstrip()

        if len(last) > max_chars_per_line - 3:
            last = last[: max_chars_per_line - 3].rstrip()

        lines[-1] = last.rstrip(" ,.-") + "..."

    return "\n".join(lines)


def _set_title(slide, text: str) -> None:
    tf = _find_text_frame(slide, "TITLE")
    if tf is None:
        return

    display_text = _prepare_slide_title(text, max_chars_per_line=32, max_lines=2)

    tf.word_wrap = True

    if not tf.paragraphs:
        tf.text = display_text
        p = tf.paragraphs[0]
    else:
        p = tf.paragraphs[0]
        p.text = display_text

    lines = display_text.splitlines() if display_text else [""]
    longest_line = max(len(line) for line in lines)

    if len(lines) == 1 and longest_line <= 20:
        size = 30
    elif len(lines) == 1 and longest_line <= 24:
        size = 28
    elif longest_line <= 24:
        size = 26
    elif longest_line <= 28:
        size = 24
    else:
        size = 22

    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    try:
        p.line_spacing = 1.0
        p.space_after = Pt(0)
    except Exception:
        pass


def _set_subtitle(slide, text: str) -> None:
    tf = _find_text_frame(slide, "SUBTITLE")
    if tf is None:
        return

    if not tf.paragraphs:
        tf.text = text
        return

    tf.paragraphs[0].text = text


def _clear_text_frame(tf) -> None:
    if tf is None:
        return

    if tf.paragraphs:
        tf.paragraphs[0].text = ""
        while len(tf.paragraphs) > 1:
            p = tf.paragraphs[-1]
            p._element.getparent().remove(p._element)
    else:
        tf.text = ""


def _split_into_bullet_lines(body_text: str) -> list[str]:
    lines: list[str] = []

    for block in (body_text or "").split("\n"):
        block = block.strip()
        if not block:
            continue

        parts = re.split(r"\.\s+", block)

        for part in parts:
            part = part.strip().rstrip(".")
            part = part.lstrip("•").lstrip("-").strip()
            if part:
                lines.append(part)

    return lines


def _set_bullets(
    slide,
    body_text: str,
    max_lines: Optional[int] = None,
    font_size_pt: int = 20,
) -> None:
    tf = _find_text_frame(slide, "BODY")
    if tf is None:
        return

    lines = _split_into_bullet_lines(body_text)

    if max_lines is not None:
        lines = lines[:max_lines]

    if not lines:
        _clear_text_frame(tf)
        return

    tf.word_wrap = True

    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]
        p._element.getparent().remove(p._element)

    p0 = tf.paragraphs[0]
    p0.text = f"• {lines[0]}"
    p0.level = 0
    p0.font.size = Pt(font_size_pt)
    p0.font.bold = False
    p0.font.name = "Segoe UI"
    p0.font.color.rgb = RGBColor(0, 0, 0)

    try:
        p0.space_after = Pt(6)
        p0.line_spacing = 1.15
    except Exception:
        pass

    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = f"• {line}"
        p.level = 0
        p.left_margin = 0
        p.first_line_indent = 0
        p.font.size = Pt(font_size_pt)
        p.font.bold = False
        p.font.name = "Segoe UI"
        p.font.color.rgb = RGBColor(0, 0, 0)

        try:
            p.space_after = Pt(6)
            p.line_spacing = 1.15
        except Exception:
            pass


# --------------------------------------------------------------------
# Speaker notes
# --------------------------------------------------------------------


def _write_speaker_notes(slide, notes_text: str) -> None:
    """
    Write speaker notes into the PowerPoint notes page.
    Safe fallback if notes are unavailable.
    """
    notes_text = (notes_text or "").strip()
    if not notes_text:
        return

    try:
        notes_slide = slide.notes_slide
    except Exception:
        return

    try:
        tf = notes_slide.notes_text_frame
        _clear_text_frame(tf)
        if tf.paragraphs:
            tf.paragraphs[0].text = notes_text
        else:
            tf.text = notes_text
        return
    except Exception:
        pass

    try:
        for shp in notes_slide.shapes:
            if not getattr(shp, "has_text_frame", False):
                continue
            try:
                ph_type = shp.placeholder_format.type
            except Exception:
                ph_type = None

            if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                tf = shp.text_frame
                _clear_text_frame(tf)
                if tf.paragraphs:
                    tf.paragraphs[0].text = notes_text
                else:
                    tf.text = notes_text
                return
    except Exception:
        return


# --------------------------------------------------------------------
# Pretty direct-render slides
# --------------------------------------------------------------------

def _render_pretty_title_slide(prs: Presentation, lesson: Dict[str, Any]) -> None:
    s = _blank_slide(prs)
    _set_slide_bg(s, RGBColor(255, 255, 255))

    # left accent bar
    bar = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.0), Inches(0.0), Inches(0.08), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    topic = lesson.get("lesson_title") or lesson.get("topic_title") or "Lesson"
    subtitle = f"{lesson.get('unit_name', '')} • {lesson.get('year_level', '')}".strip(" •")

    _add_textbox(
        s,
        Inches(0.55),
        Inches(0.55),
        Inches(8.8),
        Inches(1.1),
        text=_prepare_title_lines(topic, max_chars_per_line=26, max_lines=2),
        font_size=24,
        bold=True,
        color=TEXT,
        font_name="Aptos Display",
    )

    _add_textbox(
        s,
        Inches(0.58),
        Inches(1.42),
        Inches(6.2),
        Inches(0.4),
        text=subtitle,
        font_size=10,
        bold=False,
        color=MUTED,
        font_name="Aptos",
    )

    # Simple audience tag
    tag = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.55), Inches(6.55), Inches(1.35), Inches(0.35)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = RGBColor(245, 245, 245)
    tag.line.color.rgb = RGBColor(220, 220, 220)
    tag.line.width = Pt(1)

    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = "Audience:\n" + str(lesson.get("year_level", ""))
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT
    p.font.name = "Aptos"

    notes = lesson.get("speaker_notes_title", "") or ""
    _write_speaker_notes(s, notes)
    return s
def _render_pretty_objectives_slide(prs: Presentation, lesson: Dict[str, Any]) -> None:
    s = _blank_slide(prs)
    _set_slide_bg(s, RGBColor(255, 255, 255))

    bar = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.0), Inches(0.0), Inches(0.08), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    _add_textbox(
        s,
        Inches(0.45), Inches(0.28), Inches(3.0), Inches(0.45),
        text="Lesson Objectives",
        font_size=20,
        bold=True,
        color=TEXT,
        font_name="Aptos Display",
    )

    objectives = [o.strip() for o in lesson.get("objectives", []) if isinstance(o, str) and o.strip()]
    objectives = objectives[:3] if objectives else ["Objective 1", "Objective 2", "Objective 3"]

    _add_bullet_block(
        s,
        Inches(0.65),
        Inches(1.0),
        Inches(7.8),
        Inches(5.5),
        lines=objectives,
        font_size=18,
        color=TEXT,
        font_name="Aptos",
    )

    notes = lesson.get("speaker_notes_objectives", "") or ""
    _write_speaker_notes(s, notes)
    return s

def _render_pretty_essential_question_slide(prs: Presentation, lesson: Dict[str, Any]) -> None:
    s = _blank_slide(prs)
    _set_slide_bg(s, RGBColor(255, 255, 255))

    bar = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.0), Inches(0.0), Inches(0.08), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    _add_textbox(
        s,
        Inches(0.45), Inches(0.28), Inches(3.5), Inches(0.45),
        text="Essential Question",
        font_size=20,
        bold=True,
        color=TEXT,
        font_name="Aptos Display",
    )

    eq = (lesson.get("essential_question") or "").strip() or "What big question are we exploring today?"

    _add_textbox(
        s,
        Inches(0.8),
        Inches(1.5),
        Inches(7.2),
        Inches(1.5),
        text=eq,
        font_size=22,
        bold=False,
        color=TEXT,
        font_name="Aptos",
        align=PP_ALIGN.CENTER,
    )

    notes = lesson.get("speaker_notes_essential_question", "") or ""
    _write_speaker_notes(s, notes)
    return s


def _render_pretty_section_slide(prs: Presentation, title: str, body: str, speaker_notes: str = "") -> None:
    s = _blank_slide(prs)
    _set_slide_bg(s)
    _add_left_bar(s)
    _add_blob(s, Inches(2.75), Inches(1.15), Inches(2.2), Inches(2.2), SECONDARY, transparency=45)

    icon_map = {
        "Explore": "🔎",
        "Learn": "🧠",
        "Apply": "🛠",
        "Reflect": "💬",
    }

    clean = title.replace("🔎", "").replace("🧠", "").replace("🛠", "").replace("💬", "").strip()
    icon = icon_map.get(clean, "●")

    _add_textbox(
        s,
        Inches(3.0), Inches(1.75), Inches(1.7), Inches(0.55),
        text=f"{icon} {clean}",
        font_size=22,
        bold=True,
        color=TEXT,
        font_name="Aptos Display",
        align=PP_ALIGN.CENTER,
    )

    _add_textbox(
        s,
        Inches(2.95), Inches(2.32), Inches(1.8), Inches(0.3),
        text=f"Section transition: {clean}",
        font_size=10,
        color=MUTED,
        font_name="Aptos",
        align=PP_ALIGN.CENTER,
    )

    _write_speaker_notes(s, speaker_notes)
    return s


# --------------------------------------------------------------------
# Slide builders
# --------------------------------------------------------------------


def _add_title_slide(prs: Presentation, blueprints, lesson: Dict[str, Any]) -> None:
    return _render_pretty_title_slide(prs, lesson)


def _add_objectives_slide(prs: Presentation, blueprints, lesson: Dict[str, Any]) -> None:
    return _render_pretty_objectives_slide(prs, lesson)


def _add_essential_question_slide(prs: Presentation, blueprints, lesson: Dict[str, Any]) -> None:
    return _render_pretty_essential_question_slide(prs, lesson)


def _add_section_slide(prs: Presentation, blueprints, title: str, body: str, speaker_notes: str = "") -> None:
    return _render_pretty_section_slide(prs, title, body, speaker_notes)


def _add_visual_comparison_slide(prs: Presentation, blueprints, slide_def: Dict[str, Any]) -> None:
    """
    Create a stable two-column comparison slide using a table layout.
    """
    source = None if not blueprints else blueprints.get("content")
    s = _duplicate_slide(prs, source) if source is not None else _blank_slide(prs)

    try:
        _set_slide_bg(s)
        _add_left_bar(s)
    except Exception:
        pass

    title_text = slide_def.get("title", "") or "Comparison"
    _set_title(s, title_text)

    body_tf = _find_text_frame(s, "BODY")
    _clear_text_frame(body_tf)

    left_title = slide_def.get("left_title", "Left Side")
    left_points = slide_def.get("left_points", [])
    left_example = slide_def.get("left_example", "")

    right_title = slide_def.get("right_title", "Right Side")
    right_points = slide_def.get("right_points", [])
    right_example = slide_def.get("right_example", "")

    table_shape = s.shapes.add_table(
        rows=1,
        cols=2,
        left=Inches(0.75),
        top=Inches(1.75),
        width=Inches(8.0),
        height=Inches(4.5),
    )
    table = table_shape.table

    table.columns[0].width = Inches(3.9)
    table.columns[1].width = Inches(4.1)

    def _set_para_style(para, size_pt: int = 16, bold: bool = False) -> None:
        para.font.size = Pt(size_pt)
        para.font.bold = bold
        para.font.color.rgb = RGBColor(0, 0, 0)

    def populate_cell(cell, heading: str, points: List[str], example: str) -> None:
        cell.text = ""
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)

        tf = cell.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = heading
        _set_para_style(p, size_pt=18, bold=True)

        for point in points:
            para = tf.add_paragraph()
            para.text = f"• {point}"
            _set_para_style(para, size_pt=16, bold=False)

        if example:
            para = tf.add_paragraph()
            para.text = ""
            _set_para_style(para, size_pt=6, bold=False)

            para = tf.add_paragraph()
            para.text = "Example:"
            _set_para_style(para, size_pt=16, bold=True)

            para = tf.add_paragraph()
            para.text = example
            _set_para_style(para, size_pt=16, bold=False)

        cell.fill.background()

    populate_cell(table.cell(0, 0), left_title, left_points, left_example)
    populate_cell(table.cell(0, 1), right_title, right_points, right_example)

    line = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(4.5),
        Inches(1.65),
        Inches(0.02),
        Inches(4.65),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(210, 210, 210)
    line.line.fill.background()

    speaker_notes = slide_def.get("speaker_notes", "") or ""
    _write_speaker_notes(s, speaker_notes)

    return s


def _render_visual_comparison_slide(prs: Presentation, blueprints, slide_def: Dict[str, Any]) -> None:
    return _add_visual_comparison_slide(prs, blueprints, slide_def)


def _render_default_content_like_slide(
    prs: Presentation,
    blueprints,
    slide_def: Dict[str, Any],
) -> None:
    slide_type = (slide_def.get("type") or "content").strip()
    title_text = slide_def.get("title", "") or ""
    body_text = slide_def.get("body", "") or ""
    speaker_notes = slide_def.get("speaker_notes", "") or ""

    blueprint_map = {
        "hook": "hook",
        "video": "watch",
        "content": "content",
        "real_world": "real_world",
        "activity": "activity",
        "reflection": "reflection",
        "teacher_notes": "teacher_notes",
    }

    key = blueprint_map.get(slide_type, "content")
    source = None if not blueprints else blueprints.get(key)

    if source is not None:
        s = _duplicate_slide(prs, source)
    else:
        s = _blank_slide(prs)

    try:
        _set_slide_bg(s)
        _add_left_bar(s)
        _add_blob(s, Inches(6.9), Inches(0.12), Inches(1.0), Inches(0.42), SECONDARY, transparency=65)
        _add_blob(s, Inches(7.65), Inches(0.08), Inches(0.5), Inches(0.28), ACCENT, transparency=70)
    except Exception:
        pass

    if slide_type == "activity":
        title_text = title_text.replace(" (15 minutes)", "")
        title_text = title_text.replace("Analyse and Improve", "Analyse")
        title_text = title_text.replace("Analyze and Improve", "Analyze")

    _set_title(s, title_text)

    icon_map = {
        "hook": "💡",
        "video": "▶",
        "content": "•",
        "real_world": "🌍",
        "activity": "🛠",
        "reflection": "💬",
        "teacher_notes": "📝",
    }

    icon = icon_map.get(slide_type)

    if icon:
        try:
            _add_textbox(
                s,
                Inches(0.38),
                Inches(0.16),
                Inches(0.25),
                Inches(0.25),
                text=icon,
                font_size=13,
                bold=False,
                color=PRIMARY,
                font_name="Aptos",
            )
        except Exception:
            pass

    if slide_type == "hook":
        _set_bullets(s, body_text, max_lines=3, font_size_pt=22)
    elif slide_type == "video":
        _set_bullets(s, body_text, max_lines=3, font_size_pt=20)
    elif slide_type == "content":
        _set_bullets(s, body_text, max_lines=4, font_size_pt=20)
    elif slide_type == "real_world":
        _set_bullets(s, body_text, max_lines=4, font_size_pt=19)
    elif slide_type == "activity":
        _set_bullets(s, body_text, max_lines=6, font_size_pt=18)
    elif slide_type == "reflection":
        _set_bullets(s, body_text, max_lines=4, font_size_pt=19)
    elif slide_type == "teacher_notes":
        _set_bullets(s, body_text, max_lines=6, font_size_pt=17)
    else:
        _set_bullets(s, body_text, max_lines=4, font_size_pt=20)

    _write_speaker_notes(s, speaker_notes)

    return s


SLIDE_RENDERERS = {
    "visual_comparison": _render_visual_comparison_slide,
    "hook": _render_default_content_like_slide,
    "video": _render_default_content_like_slide,
    "content": _render_default_content_like_slide,
    "real_world": _render_default_content_like_slide,
    "activity": _render_default_content_like_slide,
    "reflection": _render_default_content_like_slide,
    "teacher_notes": _render_default_content_like_slide,
}


def _add_structured_slide(prs: Presentation, blueprints, slide_def: Dict[str, Any]) -> None:
    slide_type = (slide_def.get("type") or "content").strip()
    renderer = SLIDE_RENDERERS.get(slide_type, _render_default_content_like_slide)
    return renderer(prs, blueprints, slide_def)


# --------------------------------------------------------------------
# Public 10-slide deck builder
# --------------------------------------------------------------------

def _lesson_to_public_slide_defs(lesson: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Build a clean 10-slide public deck from a lesson JSON.
    No section slides, no teacher notes, no watch slide.
    """
    slides: List[Dict[str, Any]] = []

    source_slides = lesson.get("slides", []) or []

    def first_slide_of_type(slide_type: str) -> Optional[Dict[str, Any]]:
        for s in source_slides:
            if (s.get("type") or "").strip() == slide_type:
                return s
        return None

    def all_slides_of_type(slide_type: str) -> List[Dict[str, Any]]:
        return [s for s in source_slides if (s.get("type") or "").strip() == slide_type]

    hook_slide = first_slide_of_type("hook")
    content_slides = all_slides_of_type("content")
    real_world_slide = first_slide_of_type("real_world")
    activity_slide = first_slide_of_type("activity")
    reflection_slide = first_slide_of_type("reflection")

    slides.append({"kind": "title"})
    slides.append({"kind": "objectives"})
    slides.append({"kind": "essential_question"})

    if hook_slide:
        slides.append({
            "kind": "structured",
            "type": "hook",
            "title": "Hook Scenario",
            "body": hook_slide.get("body", ""),
            "speaker_notes": hook_slide.get("speaker_notes", ""),
        })

    for content_slide in content_slides[:3]:
        slides.append({
            "kind": "structured",
            "type": "content",
            "title": content_slide.get("title", "Key Idea"),
            "body": content_slide.get("body", ""),
            "speaker_notes": content_slide.get("speaker_notes", ""),
        })

    if real_world_slide:
        slides.append({
            "kind": "structured",
            "type": "real_world",
            "title": real_world_slide.get("title", "Real-World Example"),
            "body": real_world_slide.get("body", ""),
            "speaker_notes": real_world_slide.get("speaker_notes", ""),
        })

    if activity_slide:
        slides.append({
            "kind": "structured",
            "type": "activity",
            "title": activity_slide.get("title", "Class Activity"),
            "body": activity_slide.get("body", ""),
            "speaker_notes": activity_slide.get("speaker_notes", ""),
        })

    if reflection_slide:
        slides.append({
            "kind": "structured",
            "type": "reflection",
            "title": reflection_slide.get("title", "Reflection Questions"),
            "body": reflection_slide.get("body", ""),
            "speaker_notes": reflection_slide.get("speaker_notes", ""),
        })

    return slides[:10]


# --------------------------------------------------------------------
# Main conversion
# --------------------------------------------------------------------


def lesson_json_to_pptx(lesson_path: Path, output_dir: Path | None = None) -> Path:
    lesson = _load_lesson(lesson_path)

    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE_PATH}")

    prs = Presentation(str(TEMPLATE_PATH))

    original_rids = {sldId.rId for sldId in list(prs.slides._sldIdLst)}  # type: ignore[attr-defined]
    blueprints = _index_blueprints(prs)

    _add_title_slide(prs, blueprints, lesson)
    _add_objectives_slide(prs, blueprints, lesson)
    _add_essential_question_slide(prs, blueprints, lesson)

    for sd in lesson.get("slides", []):
        if not isinstance(sd, dict):
            continue

        if (sd.get("type") or "").strip() == "section":
            _add_section_slide(
                prs,
                blueprints,
                sd.get("title", ""),
                sd.get("body", ""),
                sd.get("speaker_notes", "") or "",
            )
        else:
            _add_structured_slide(prs, blueprints, sd)

    _delete_slides_by_rids(prs, original_rids)

    if output_dir is None:
        output_dir = Path("output")
    output_dir.mkdir(parents=True, exist_ok=True)

    topic_slug = slugify(lesson.get("topic_title", lesson.get("lesson_title", "lesson")))
    lesson_number = lesson.get("lesson_number")
    prefix = f"{lesson_number:02d}-" if isinstance(lesson_number, int) else ""
    output_path = output_dir / f"{prefix}{topic_slug}.pptx"

    prs.save(output_path)
    return output_path

# --------------------------------------------------------------------
# Public deck direct-render slides
# --------------------------------------------------------------------
def _render_public_structured_slide(prs: Presentation, slide_def: Dict[str, Any]) -> None:
    """
    Simple, clean public slide renderer.
    Minimalist layout designed for readability and ALai post-processing.
    """
    s = _blank_slide(prs)

    # Clean white background
    _set_slide_bg(s, RGBColor(255, 255, 255))

    # Thin left accent bar
    bar = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.0), Inches(0.0), Inches(0.08), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    slide_type = (slide_def.get("type") or "content").strip()
    title_text = (slide_def.get("title") or "").strip()
    body_text = (slide_def.get("body") or "").strip()
    speaker_notes = (slide_def.get("speaker_notes") or "").strip()

    # Title
    _add_textbox(
        s,
        Inches(0.45),
        Inches(0.28),
        Inches(8.6),
        Inches(0.55),
        text=_prepare_title_lines(title_text, max_chars_per_line=34, max_lines=2),
        font_size=20,
        bold=True,
        color=TEXT,
        font_name="Aptos Display",
    )

    # Body bullets
    lines = _split_into_bullet_lines(body_text)

    if slide_type == "hook":
        max_lines = 4
        font_size = 20
    elif slide_type == "activity":
        max_lines = 8
        font_size = 16
    elif slide_type == "reflection":
        max_lines = 5
        font_size = 18
    else:
        max_lines = 5
        font_size = 18

    lines = lines[:max_lines]

    _add_bullet_block(
        s,
        Inches(0.6),
        Inches(1.05),
        Inches(8.0),
        Inches(5.9),
        lines=lines,
        font_size=font_size,
        color=TEXT,
        font_name="Aptos",
    )

    _write_speaker_notes(s, speaker_notes)
    return s

def lesson_json_to_public_pptx(
    lesson_json_path: Path,
    output_dir: Path | None = None,
) -> Path:
    """
    Generate a concise 10-slide public-facing PPTX deck from a lesson JSON.
    Intended for student/public presentation and ALai post-processing.
    """
    lesson = _load_lesson(lesson_json_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    public_slide_defs = _lesson_to_public_slide_defs(lesson)

    for slide_def in public_slide_defs:
        kind = slide_def.get("kind")

        if kind == "title":
            _render_pretty_title_slide(prs, lesson)
        elif kind == "objectives":
            _render_pretty_objectives_slide(prs, lesson)
        elif kind == "essential_question":
            _render_pretty_essential_question_slide(prs, lesson)
        elif kind == "structured":
            _render_public_structured_slide(prs, slide_def)

    if output_dir is None:
        output_dir = lesson_json_path.parent
    output_dir.mkdir(parents=True, exist_ok=True)

    topic_slug = slugify(lesson.get("topic_title", lesson.get("lesson_title", "lesson")))
    lesson_number = lesson.get("lesson_number")
    prefix = f"{lesson_number:02d}-" if isinstance(lesson_number, int) else ""

    out_path = output_dir / f"{prefix}{topic_slug}-public.pptx"
    prs.save(out_path)

    return out_path


if __name__ == "__main__":
    sample = next(BASE_OUTPUT_DIR.rglob("lesson.json"))
    out = lesson_json_to_pptx(sample)
    print(f"Saved PPTX to {out}")