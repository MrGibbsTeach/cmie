"""
Premium PPTX generator for CMIE lesson slides.
Produces visually polished, market-ready presentations.
"""
from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .ai_lesson_engine import slugify

SW = Inches(13.333)
SH = Inches(7.5)

# ── Colour palette ────────────────────────────────────────────────
INDIGO      = RGBColor(67,  56,  202)
INDIGO_MED  = RGBColor(99,  102, 241)
INDIGO_DIM  = RGBColor(55,  48,  163)
INDIGO_SOFT = RGBColor(199, 210, 254)
INDIGO_BG   = RGBColor(238, 242, 255)
TEAL        = RGBColor(13,  148, 136)
TEAL_BG     = RGBColor(240, 253, 250)
TEAL_SOFT   = RGBColor(153, 246, 228)
AMBER       = RGBColor(180, 83,   9)
AMBER_BG    = RGBColor(255, 251, 235)
AMBER_SOFT  = RGBColor(253, 230, 138)
ROSE        = RGBColor(190, 18,  60)
NAVY        = RGBColor(15,  23,  42)
SLATE       = RGBColor(100, 116, 139)
SLATE_LIGHT = RGBColor(226, 232, 240)
WHITE       = RGBColor(255, 255, 255)
OFF_WHITE   = RGBColor(248, 250, 252)
WARM_WHITE  = RGBColor(255, 252, 245)

FH = "Aptos Display"
FB = "Aptos"


# ── Shape helpers ─────────────────────────────────────────────────

def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def _rect(slide, l, t, w, h, fill, line=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _rrect(slide, l, t, w, h, fill, line=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _oval(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1.0)
    else:
        s.line.fill.background()
    return s

def _txt(slide, l, t, w, h, text, size, bold=False, color=NAVY,
         align=PP_ALIGN.LEFT, font=FB, wrap=True, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box

def _bullets(slide, l, t, w, h, lines, size, color=NAVY, font=FB,
             spacing=12, line_spacing=1.25, marker="•"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{marker}  {line}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font
        try:
            p.space_after = Pt(spacing)
            p.line_spacing = line_spacing
        except Exception:
            pass
    return box

def _hyperlink_txt(slide, l, t, w, h, display, url, size,
                   color=INDIGO_MED, font=FB):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = display
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font
    run.font.underline = True
    try:
        run.hyperlink.address = url
    except Exception:
        pass
    return box

def _split_lines(body, max_lines=None):
    lines = []
    for raw in (body or "").splitlines():
        raw = raw.strip().lstrip("•").lstrip("-").strip()
        if not raw:
            continue
        if len(raw) > 120:
            for part in re.split(r"\.\s+", raw):
                part = part.strip().rstrip(".")
                if part:
                    lines.append(part)
        else:
            lines.append(raw)
    return lines[:max_lines] if max_lines else lines

def _header(slide, title, band_color, band_h=1.4, prefix=""):
    _rect(slide, Inches(0), Inches(0), SW, Inches(band_h), fill=band_color)
    display = f"{prefix}{title}" if prefix else title
    _txt(slide, Inches(0.55), Inches(0.16), Inches(12.3), Inches(band_h - 0.08),
         display, size=28, bold=True, color=WHITE, font=FH)

def _badge(slide, l, t, w, h, text, bg, fg=WHITE, size=12):
    s = _rrect(slide, l, t, w, h, fill=bg, line=None)
    tf = s.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = fg
    p.font.name = FB
    p.alignment = PP_ALIGN.CENTER
    return s


# ── Slide renderers ───────────────────────────────────────────────

def _render_title(prs, lesson):
    s = _slide(prs)
    _bg(s, WHITE)

    panel_w = Inches(7.5)

    # Left INDIGO panel
    _rect(s, Inches(0), Inches(0), panel_w, SH, fill=INDIGO)
    # Thin dark accent strip at the seam
    _rect(s, panel_w - Inches(0.07), Inches(0), Inches(0.14), SH, fill=INDIGO_DIM)

    # Left panel: large translucent lesson number as background element
    num = lesson.get("lesson_number")
    num_str = (f"{num:02d}" if isinstance(num, int) else str(num)) if num else ""
    if num_str:
        _txt(s, Inches(-0.4), Inches(2.8), Inches(5.0), Inches(5.5),
             num_str, size=210, bold=True, color=INDIGO_DIM, font=FH,
             align=PP_ALIGN.CENTER)

    # Left panel: subject label (small, top)
    subject = lesson.get("subject", "")
    if subject:
        _txt(s, Inches(0.55), Inches(0.38), Inches(6.5), Inches(0.45),
             subject.upper(), size=11, color=INDIGO_SOFT, font=FB)

    # Left panel: amber accent rule
    _rect(s, Inches(0.55), Inches(1.35), Inches(2.6), Inches(0.08), fill=AMBER_SOFT)

    # Left panel: title (dominant)
    title = lesson.get("lesson_title") or lesson.get("topic_title") or "Lesson"
    _txt(s, Inches(0.55), Inches(1.62), Inches(6.65), Inches(3.4),
         title, size=42, bold=True, color=WHITE, font=FH)

    # Left panel: unit name
    unit = lesson.get("unit_name", "")
    if unit:
        _txt(s, Inches(0.55), Inches(5.0), Inches(6.5), Inches(0.55),
             unit, size=14, color=INDIGO_SOFT, font=FB)

    # Left panel: lesson badge
    if num_str:
        _badge(s, Inches(0.55), Inches(6.62), Inches(1.92), Inches(0.52),
               f"Lesson {num_str}", WHITE, INDIGO, 12)

    # Right OFF_WHITE panel
    _rect(s, panel_w, Inches(0), SW - panel_w, SH, fill=OFF_WHITE)

    # Right panel: decorative circle composition
    _oval(s, Inches(9.0), Inches(0.0), Inches(3.5), Inches(3.5), INDIGO_BG)
    _oval(s, Inches(11.4), Inches(-0.4), Inches(2.3), Inches(2.3), INDIGO_SOFT)
    _oval(s, Inches(8.5), Inches(4.2), Inches(2.8), Inches(2.8), INDIGO_BG)
    _oval(s, Inches(11.8), Inches(5.2), Inches(1.7), Inches(1.7), INDIGO_SOFT)

    # Right panel: subject / year badge
    year = lesson.get("year_level", "")
    badge_text = "  ·  ".join(x for x in [subject, year] if x)
    if badge_text:
        b = _rrect(s, panel_w + Inches(0.28), Inches(3.55), Inches(5.28), Inches(0.5),
                   fill=WHITE, line=SLATE_LIGHT, lw=1.0)
        tf = b.text_frame
        p = tf.paragraphs[0]
        p.text = badge_text
        p.font.size = Pt(12)
        p.font.color.rgb = SLATE
        p.font.name = FB
        p.alignment = PP_ALIGN.CENTER

    _write_notes(s, lesson.get("speaker_notes_title", "") or "")


def _render_objectives(prs, lesson):
    s = _slide(prs)
    _bg(s, WHITE)

    # Header with subtitle
    _rect(s, Inches(0), Inches(0), SW, Inches(1.6), fill=INDIGO)
    _txt(s, Inches(0.55), Inches(0.12), Inches(10), Inches(0.88),
         "Learning Objectives", size=30, bold=True, color=WHITE, font=FH)
    _txt(s, Inches(0.55), Inches(0.98), Inches(10), Inches(0.48),
         "By the end of this lesson, you will be able to:", size=13,
         color=INDIGO_SOFT, font=FB, italic=True)

    objectives = [o.strip() for o in lesson.get("objectives", [])
                  if isinstance(o, str) and o.strip()][:5]

    card_h = Inches(1.0)
    gap = Inches(0.1)
    top_start = Inches(1.72)
    bc_list = [INDIGO, TEAL, AMBER, ROSE, INDIGO_MED]

    for i, obj in enumerate(objectives):
        top = top_start + (card_h + gap) * i
        bc = bc_list[i % len(bc_list)]

        # Shadow strip for depth
        _rrect(s, Inches(0.5), top + Inches(0.06), Inches(12.43), card_h,
               fill=SLATE_LIGHT, line=None)
        # Card body
        _rrect(s, Inches(0.45), top, Inches(12.43), card_h,
               fill=OFF_WHITE, line=SLATE_LIGHT, lw=0.5)
        # Colored left accent
        _rect(s, Inches(0.45), top, Inches(0.12), card_h, fill=bc)
        # Number
        _txt(s, Inches(0.7), top + Inches(0.2), Inches(0.52), Inches(0.58),
             str(i + 1), size=26, bold=True, color=bc, font=FH)
        # Objective text
        _txt(s, Inches(1.42), top + Inches(0.15), Inches(11.1), Inches(0.78),
             obj, size=20, color=NAVY, font=FB)

    _write_notes(s, lesson.get("speaker_notes_objectives", "") or "")


def _render_essential_question(prs, lesson):
    s = _slide(prs)
    # Full dark background — dramatic and premium
    _bg(s, NAVY)

    # Large watermark open-quote (top right, INDIGO_DIM barely visible on dark)
    _txt(s, Inches(7.5), Inches(-2.0), Inches(6.5), Inches(6.0),
         "“", size=380, bold=True, color=INDIGO_DIM, font=FH,
         align=PP_ALIGN.CENTER)

    # Label badge
    _badge(s, Inches(0.55), Inches(0.42), Inches(2.5), Inches(0.5),
           "Essential Question", INDIGO, WHITE, 12)

    # Thin accent line
    _rect(s, Inches(0.55), Inches(1.08), Inches(2.5), Inches(0.04), fill=INDIGO_MED)

    eq = (lesson.get("essential_question") or "What question are we exploring today?").strip()

    # Opening quote (visible)
    _txt(s, Inches(0.5), Inches(1.15), Inches(1.4), Inches(1.4),
         "“", size=110, bold=True, color=INDIGO_SOFT, font=FH)

    # The question — large, white, confident
    _txt(s, Inches(0.65), Inches(2.3), Inches(10.8), Inches(4.2),
         eq, size=34, bold=False, color=WHITE, font=FH)

    # Closing quote (bottom right, subtle)
    _txt(s, Inches(10.0), Inches(5.0), Inches(3.0), Inches(2.0),
         "”", size=110, bold=True, color=INDIGO_DIM, font=FH)

    _write_notes(s, lesson.get("speaker_notes_essential_question", "") or "")


def _render_section(prs, slide_def):
    s = _slide(prs)
    _bg(s, INDIGO)

    # Background blobs (organic shapes)
    _oval(s, Inches(-1.2), Inches(3.8), Inches(5.5), Inches(5.5), INDIGO_MED)
    _oval(s, Inches(9.8), Inches(-1.5), Inches(5.0), Inches(5.0), INDIGO_MED)
    _oval(s, Inches(10.3), Inches(4.8), Inches(4.0), Inches(4.0), INDIGO_DIM)

    raw_title = slide_def.get("title", "")
    clean = re.sub(r"[🔎🧠🛠💬🔍]", "", raw_title).strip()
    icon_map = {"Explore": "🔎", "Learn": "🧠", "Apply": "🛠", "Reflect": "💬"}
    icon = icon_map.get(clean, "●")

    # Subtle centred rule
    _rect(s, Inches(3.8), Inches(3.52), Inches(5.7), Inches(0.05), fill=INDIGO_SOFT)

    # Icon
    _txt(s, Inches(0), Inches(1.4), SW, Inches(2.0),
         icon, size=82, color=WHITE, font=FB, align=PP_ALIGN.CENTER)

    # Section name
    _txt(s, Inches(0), Inches(3.6), SW, Inches(1.6),
         clean.upper(), size=52, bold=True, color=WHITE, font=FH,
         align=PP_ALIGN.CENTER)

    body = slide_def.get("body", "")
    if body:
        _txt(s, Inches(0), Inches(5.28), SW, Inches(0.9),
             body, size=16, color=INDIGO_SOFT, font=FB, align=PP_ALIGN.CENTER)

    _write_notes(s, slide_def.get("speaker_notes", "") or "")


def _render_hook(prs, slide_def):
    s = _slide(prs)
    _bg(s, INDIGO_BG)

    # Badge
    _badge(s, Inches(0.55), Inches(0.32), Inches(1.62), Inches(0.5),
           "💡 Hook", INDIGO, WHITE, 13)

    # Decorative quote mark (background, top right)
    _txt(s, Inches(10.0), Inches(0.0), Inches(3.5), Inches(4.2),
         "“", size=260, bold=True, color=INDIGO_SOFT, font=FH)

    # White card with left accent
    _rrect(s, Inches(0.55), Inches(1.02), Inches(12.25), Inches(5.85),
           fill=WHITE, line=INDIGO_SOFT, lw=1.5)
    _rect(s, Inches(0.55), Inches(1.02), Inches(0.12), Inches(5.85), fill=INDIGO)

    body = slide_def.get("body", "")
    lines = _split_lines(body, max_lines=5)
    if lines:
        _txt(s, Inches(1.05), Inches(1.6), Inches(11.35), Inches(5.0),
             "\n\n".join(lines), size=22, color=NAVY, font=FB)

    _write_notes(s, slide_def.get("speaker_notes", "") or "")


def _render_content(prs, slide_def):
    s = _slide(prs)
    # Floating card: light lavender bg + white card = premium depth
    _bg(s, INDIGO_BG)
    _header(s, slide_def.get("title", ""), INDIGO, band_h=1.42)

    # Shadow for depth
    _rrect(s, Inches(0.38), Inches(1.56), Inches(12.73), Inches(5.66),
           fill=SLATE_LIGHT, line=None)
    # White floating card
    _rrect(s, Inches(0.3), Inches(1.48), Inches(12.73), Inches(5.66),
           fill=WHITE, line=None)
    # Left accent bar on card
    _rect(s, Inches(0.3), Inches(1.48), Inches(0.14), Inches(5.66), fill=INDIGO)

    lines = _split_lines(slide_def.get("body", ""), max_lines=5)
    _bullets(s, Inches(0.64), Inches(1.82), Inches(12.1), Inches(5.0),
             lines, size=20, color=NAVY)

    _write_notes(s, slide_def.get("speaker_notes", "") or "")


def _render_real_world(prs, slide_def):
    s = _slide(prs)
    # Floating card with teal theme
    _bg(s, TEAL_BG)

    title = slide_def.get("title", "Real-World Example")
    clean_title = re.sub(r"[🌍]", "", title).strip()
    _header(s, clean_title, TEAL, band_h=1.42, prefix="🌍  ")

    # Shadow + white floating card
    _rrect(s, Inches(0.38), Inches(1.56), Inches(12.73), Inches(5.66),
           fill=TEAL_SOFT, line=None)
    _rrect(s, Inches(0.3), Inches(1.48), Inches(12.73), Inches(5.66),
           fill=WHITE, line=None)
    _rect(s, Inches(0.3), Inches(1.48), Inches(0.14), Inches(5.66), fill=TEAL)

    lines = _split_lines(slide_def.get("body", ""), max_lines=5)
    _bullets(s, Inches(0.64), Inches(1.82), Inches(12.1), Inches(5.0),
             lines, size=20, color=NAVY, marker="→")

    _write_notes(s, slide_def.get("speaker_notes", "") or "")


def _render_activity(prs, slide_def):
    s = _slide(prs)
    _bg(s, AMBER_BG)

    title = re.sub(r"\(\d+ minutes?\)", "", slide_def.get("title", "Activity")).strip()
    _header(s, title, AMBER, band_h=1.42)

    # Time badge
    raw_title = slide_def.get("title", "")
    m = re.search(r"\((\d+\s*minutes?)\)", raw_title, re.IGNORECASE)
    if m:
        _badge(s, Inches(11.0), Inches(0.34), Inches(2.0), Inches(0.5),
               f"⏱  {m.group(1)}", WHITE, AMBER, 12)

    # Split body into steps and output section
    lines = _split_lines(slide_def.get("body", ""), max_lines=10)
    step_lines, output_lines = [], []
    in_output = False
    for line in lines:
        lo = line.lower()
        if lo.startswith("output") or lo.startswith("you must"):
            in_output = True
        if in_output:
            output_lines.append(line)
        else:
            step_lines.append(line)

    # Numbered amber step circles
    max_steps = 4 if output_lines else 5
    for i, line in enumerate(step_lines[:max_steps]):
        top = Inches(1.6) + Inches(i * 0.96)
        circle = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,
                                     Inches(0.45), top, Inches(0.52), Inches(0.52))
        circle.fill.solid()
        circle.fill.fore_color.rgb = AMBER
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FB
        p.alignment = PP_ALIGN.CENTER
        _txt(s, Inches(1.12), top + Inches(0.04), Inches(11.9), Inches(0.58),
             line, size=18, color=NAVY, font=FB)

    # Output box
    if output_lines:
        _rrect(s, Inches(0.3), Inches(5.35), Inches(12.73), Inches(1.88),
               fill=AMBER_SOFT, line=AMBER, lw=1.5)
        # Output label pill
        _badge(s, Inches(0.5), Inches(5.42), Inches(1.62), Inches(0.45),
               "📋 Output", AMBER, WHITE, 11)
        # Clean the output text
        cleaned = []
        for ol in output_lines:
            c = re.sub(r'^(output|you must)\s*:?\s*', '', ol, flags=re.IGNORECASE).strip()
            if c and c not in ("|", ""):
                cleaned.append(c)
        if cleaned:
            _txt(s, Inches(2.3), Inches(5.5), Inches(10.6), Inches(0.65),
                 cleaned[0], size=14, color=NAVY, font=FB)
        if len(cleaned) > 1:
            _txt(s, Inches(0.5), Inches(6.15), Inches(12.4), Inches(0.55),
                 f"✓  {cleaned[1]}", size=13, bold=True, color=AMBER, font=FB)

    _write_notes(s, slide_def.get("speaker_notes", "") or "")


def _render_reflection(prs, slide_def):
    s = _slide(prs)
    _bg(s, WHITE)

    # INDIGO_BG header with icon
    _rect(s, Inches(0), Inches(0), SW, Inches(1.52), fill=INDIGO_BG)
    _txt(s, Inches(0.55), Inches(0.22), Inches(12.3), Inches(1.08),
         "💬  " + slide_def.get("title", "Reflection"), size=28,
         bold=True, color=INDIGO, font=FH)

    lines = _split_lines(slide_def.get("body", ""), max_lines=4)
    card_h = Inches(1.32)
    gap = Inches(0.1)
    top_start = Inches(1.65)

    for i, line in enumerate(lines[:4]):
        top = top_start + (card_h + gap) * i
        # Shadow
        _rrect(s, Inches(0.5), top + Inches(0.07), Inches(12.43), card_h,
               fill=SLATE_LIGHT, line=None)
        # Card
        _rrect(s, Inches(0.45), top, Inches(12.43), card_h,
               fill=INDIGO_BG, line=INDIGO_SOFT, lw=1.0)
        # Left accent
        _rect(s, Inches(0.45), top, Inches(0.11), card_h, fill=INDIGO_MED)
        # Text
        _txt(s, Inches(0.76), top + Inches(0.24),
             Inches(11.9), Inches(0.9), line, size=20, color=NAVY, font=FB)

    _write_notes(s, slide_def.get("speaker_notes", "") or "")


def _render_teacher_notes(prs, slide_def):
    s = _slide(prs)
    _bg(s, WARM_WHITE)

    _rect(s, Inches(0), Inches(0), SW, Inches(1.38), fill=AMBER)
    _txt(s, Inches(0.55), Inches(0.15), Inches(12), Inches(1.1),
         "📝  Teacher Notes", size=24, bold=True, color=WHITE, font=FH)

    # Watermark
    _txt(s, Inches(0), Inches(2.5), SW, Inches(2.8),
         "TEACHER", size=96, bold=True,
         color=AMBER_SOFT, font=FH, align=PP_ALIGN.CENTER)

    lines = _split_lines(slide_def.get("body", ""), max_lines=8)
    _bullets(s, Inches(0.55), Inches(1.52), Inches(12.3), Inches(5.6),
             lines, size=17, color=NAVY, spacing=8)

    _write_notes(s, slide_def.get("speaker_notes", "") or "")


def _render_video(prs, slide_def, video_url=None):
    s = _slide(prs)
    _bg(s, WHITE)
    _header(s, slide_def.get("title", "Watch & Think"), INDIGO_MED, band_h=1.42)

    vl = Inches(0.42)
    vt = Inches(1.55)
    vw = Inches(7.58)
    vh = Inches(4.9)
    bar_h = Inches(0.45)

    # ── Main video box (clickable if URL present) ─────────────────
    video_box = _rrect(s, vl, vt, vw, vh, fill=NAVY, line=None)
    if video_url:
        try:
            video_box.click_action.hyperlink.address = video_url
        except Exception:
            pass

    # Subtle inner glow top edge (lighter strip at top of video box)
    _rect(s, vl, vt, vw, Inches(0.08), fill=RGBColor(30, 40, 70))

    # YouTube-style red control bar (bottom of video box)
    _rect(s, vl, vt + vh - bar_h, vw, bar_h, fill=RGBColor(220, 38, 38))

    # Progress indicator (white thin bar inside the red bar)
    _rect(s, vl + Inches(0.18), vt + vh - bar_h + Inches(0.1),
          vw * 0.55, Inches(0.06), fill=WHITE)

    # Play button circle (red, matching YouTube)
    pr = Inches(0.85)
    pcx = vl + vw / 2
    pcy = vt + (vh - bar_h) / 2
    play_circle = _oval(s, pcx - pr / 2, pcy - pr / 2, pr, pr,
                        RGBColor(220, 38, 38))
    tf = play_circle.text_frame
    p = tf.paragraphs[0]
    p.text = "▶"
    p.font.size = Pt(30)
    p.font.color.rgb = WHITE
    p.font.name = FB
    p.alignment = PP_ALIGN.CENTER

    # "CLICK TO WATCH" overlay (top of video box, only when URL exists)
    if video_url:
        _txt(s, vl, vt + Inches(0.18), vw, Inches(0.48),
             "▶  CLICK TO WATCH IN PRESENTATION MODE",
             size=10, bold=True, color=RGBColor(160, 170, 210),
             font=FB, align=PP_ALIGN.CENTER)

    # URL hyperlink below the box
    if video_url:
        _hyperlink_txt(s, vl, vt + vh + Inches(0.14), vw, Inches(0.44),
                       video_url, video_url, size=11)
    else:
        _txt(s, vl, vt + vh + Inches(0.14), vw, Inches(0.44),
             "Video URL will be provided by your teacher",
             size=11, color=SLATE, font=FB, italic=True)

    # ── Right side: watch prompt cards ───────────────────────────
    lines = _split_lines(slide_def.get("body", ""), max_lines=4)
    prompts = [l for l in lines
               if "teacher" not in l.lower() and "show" not in l.lower()][:3]

    if prompts:
        _txt(s, Inches(8.28), Inches(1.72), Inches(4.75), Inches(0.5),
             "While you watch:", size=14, bold=True, color=SLATE, font=FB)
        for j, prompt in enumerate(prompts):
            ct = Inches(2.38) + Inches(j * 1.52)
            # Shadow
            _rrect(s, Inches(8.35), ct + Inches(0.07), Inches(4.75), Inches(1.32),
                   fill=SLATE_LIGHT, line=None)
            # Card
            _rrect(s, Inches(8.28), ct, Inches(4.75), Inches(1.32),
                   fill=INDIGO_BG, line=INDIGO_SOFT, lw=0.8)
            _rect(s, Inches(8.28), ct, Inches(0.08), Inches(1.32), fill=INDIGO_MED)
            _txt(s, Inches(8.56), ct + Inches(0.2),
                 Inches(4.28), Inches(0.9), prompt, size=16, color=NAVY, font=FB)

    _write_notes(s, slide_def.get("speaker_notes", "") or "")


# ── Speaker notes ─────────────────────────────────────────────────

def _write_notes(slide, text):
    if not text:
        return
    try:
        tf = slide.notes_slide.notes_text_frame
        if tf.paragraphs:
            tf.paragraphs[0].text = text
        else:
            tf.text = text
    except Exception:
        pass


# ── Renderer dispatch ─────────────────────────────────────────────

_RENDERERS = {
    "section":       _render_section,
    "hook":          _render_hook,
    "video":         _render_video,
    "content":       _render_content,
    "real_world":    _render_real_world,
    "activity":      _render_activity,
    "reflection":    _render_reflection,
    "teacher_notes": _render_teacher_notes,
}


# ── Public entry point ────────────────────────────────────────────

def lesson_json_to_pptx(lesson_path: Path, output_dir: Optional[Path] = None) -> Path:
    """Convert a lesson JSON file to a premium PPTX presentation."""
    with lesson_path.open(encoding="utf-8") as f:
        lesson: Dict[str, Any] = json.load(f)

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    _render_title(prs, lesson)
    _render_objectives(prs, lesson)
    _render_essential_question(prs, lesson)

    video_url = lesson.get("video_url")

    for sd in lesson.get("slides", []):
        if not isinstance(sd, dict):
            continue
        stype = (sd.get("type") or "content").strip()
        renderer = _RENDERERS.get(stype, _render_content)
        if stype == "video":
            renderer(prs, sd, video_url)
        else:
            renderer(prs, sd)

    if output_dir is None:
        output_dir = lesson_path.parent
    output_dir.mkdir(parents=True, exist_ok=True)

    topic_slug = slugify(lesson.get("topic_title") or lesson.get("lesson_title") or "lesson")
    num = lesson.get("lesson_number")
    prefix = f"{num:02d}-" if isinstance(num, int) else ""
    out = output_dir / f"{prefix}{topic_slug}.pptx"
    prs.save(str(out))
    return out
