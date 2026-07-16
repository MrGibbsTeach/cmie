"""
make_lead_magnet.py — turn a unit's Lesson 1 into a standalone free-resource
package: an exact copy of the Lesson 1 deck with one extra slide appended
("enjoyed this? here's the full unit, and a review would mean a lot").

This does NOT touch the existing paid Lesson 1 listing or reprice anything —
it produces a new, separate zip. Publishing it (as a free product) is a
manual/separate step, same as every other artifact this pipeline builds.

Usage:
    python make_lead_magnet.py --unit year7_networks_hardware_unit1
    python make_lead_magnet.py --unit year7_networks_hardware_unit1 \
        --bundle-url https://www.teacherspayteachers.com/Product/...
"""
from __future__ import annotations

import argparse
import shutil
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).parent
PUBLIC_ROOT = PROJECT_ROOT / "releases" / "public"
ARTIFACTS_ROOT = PROJECT_ROOT / "releases" / "artifacts"

SW = Inches(13.333)
SH = Inches(7.5)

INDIGO = RGBColor(67, 56, 202)
INDIGO_BG = RGBColor(238, 242, 255)
TEAL = RGBColor(13, 148, 136)
NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
FH = "Aptos Display"
FB = "Aptos"

STORE_NAME = "FocusLab Digital"


def _find_lesson01_pptx(unit_id: str, version: str) -> Path:
    slides_dir = PUBLIC_ROOT / f"{unit_id}_{version}" / "01_Lesson_Slides"
    if not slides_dir.exists():
        raise FileNotFoundError(f"No slides folder for {unit_id}: {slides_dir}")
    matches = sorted(slides_dir.glob("01-*.pptx"))
    if not matches:
        raise FileNotFoundError(f"No Lesson 1 pptx found in {slides_dir}")
    return matches[0]


def _read_unit_title(unit_id: str, version: str) -> str:
    listing = PUBLIC_ROOT / f"{unit_id}_{version}" / "06_Listings" / "unit" / "tpt_listing.md"
    if listing.exists():
        first_line = listing.read_text(encoding="utf-8").splitlines()[0]
        return first_line.lstrip("#").strip()
    return unit_id.replace("_", " ").title()


def _rect(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # 1 = MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _rrect(slide, l, t, w, h, fill):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _txt(slide, l, t, w, h, text, size, bold=False, color=NAVY,
         align=PP_ALIGN.LEFT, font=FB, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box


def _hyperlink_txt(slide, l, t, w, h, display, url, size, color=INDIGO):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = display
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = FB
    run.font.underline = True
    try:
        run.hyperlink.address = url
    except Exception:
        pass
    return box


def _add_cta_slide(prs: Presentation, unit_title: str, bundle_url: str | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    _rect(slide, Inches(0), Inches(0), SW, SH, fill=WHITE)
    _rect(slide, Inches(0), Inches(0), SW, Inches(1.4), fill=INDIGO)
    _txt(slide, Inches(0.55), Inches(0.16), Inches(12.3), Inches(1.2),
         "Thanks for teaching this lesson!", size=28, bold=True, color=WHITE, font=FH)

    _rrect(slide, Inches(1.4), Inches(1.9), Inches(10.5), Inches(2.6), fill=INDIGO_BG)
    _txt(slide, Inches(1.9), Inches(2.15), Inches(9.5), Inches(0.6),
         "Want the full unit?", size=22, bold=True, color=INDIGO, font=FH)
    _txt(slide, Inches(1.9), Inches(2.75), Inches(9.5), Inches(1.5),
         f'"{unit_title}" includes all 7 lessons, a student workbook, unit roadmap, '
         "and a full assessment pack with rubric — everything you need to teach "
         "the whole unit, not just this lesson.",
         size=15, color=NAVY, font=FB)

    if bundle_url:
        _hyperlink_txt(slide, Inches(2.4), Inches(4.55), Inches(8.5), Inches(0.5),
                        "Get the full unit →", bundle_url, size=16)
    else:
        _txt(slide, Inches(1.9), Inches(4.55), Inches(9.5), Inches(0.5),
             f"Search “{unit_title}” by {STORE_NAME}", size=16,
             bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

    _txt(slide, Inches(1.4), Inches(5.5), Inches(10.5), Inches(0.5),
         "Enjoying this resource?", size=18, bold=True, color=TEAL,
         align=PP_ALIGN.CENTER, font=FH)
    _txt(slide, Inches(1.9), Inches(6.0), Inches(9.5), Inches(0.8),
         "A quick review helps a small independent teacher-store like ours more "
         "than you'd think — thank you for the support!",
         size=14, color=SLATE, align=PP_ALIGN.CENTER, font=FB)
    _txt(slide, Inches(1.4), Inches(7.05), Inches(10.5), Inches(0.35),
         STORE_NAME, size=11, color=SLATE, align=PP_ALIGN.CENTER, font=FB)


def make_lead_magnet(unit_id: str, version: str = "v001", bundle_url: str | None = None,
                      out_dir: Path | None = None) -> Path:
    lesson01_path = _find_lesson01_pptx(unit_id, version)
    unit_title = _read_unit_title(unit_id, version)

    out_dir = out_dir or ARTIFACTS_ROOT
    out_dir.mkdir(parents=True, exist_ok=True)

    work_pptx = out_dir / f"{lesson01_path.stem}_FREE_SAMPLE.pptx"
    shutil.copy2(lesson01_path, work_pptx)

    prs = Presentation(str(work_pptx))
    prs.slide_width = SW
    prs.slide_height = SH
    _add_cta_slide(prs, unit_title, bundle_url)
    prs.save(str(work_pptx))

    zip_path = out_dir / f"{unit_id}_lesson01_FREE_{version}.zip"
    if zip_path.exists():
        zip_path.unlink()
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.write(work_pptx, arcname=work_pptx.name)
    work_pptx.unlink()

    with zipfile.ZipFile(zip_path) as zf:
        bad = zf.testzip()
        if bad:
            raise RuntimeError(f"Corrupt entry in {zip_path}: {bad}")
        names = zf.namelist()

    print(f"Lead magnet built: {zip_path}")
    print(f"  Contains: {names}")
    print(f"  Unit title used in CTA: {unit_title}")
    print(f"  Bundle link: {bundle_url or '(generic search text, no URL provided)'}")
    return zip_path


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--unit", required=True, help="Unit ID, e.g. year7_networks_hardware_unit1")
    parser.add_argument("--version", default="v001")
    parser.add_argument("--bundle-url", default=None,
                         help="Real TPT/TES bundle product URL to link to (optional; "
                              "falls back to a generic 'search for this title' CTA)")
    parser.add_argument("--out-dir", default=None)
    args = parser.parse_args()

    out_dir = Path(args.out_dir) if args.out_dir else None
    try:
        make_lead_magnet(args.unit, args.version, args.bundle_url, out_dir)
    except Exception as e:
        print(f"ERROR: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
